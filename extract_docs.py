#!/usr/bin/env python3
"""
Script para extraer contenido de documentos Word, PowerPoint y Excel
"""

import os
import sys
from pathlib import Path

def extract_docx_content(file_path):
    """Extraer contenido de archivos .docx"""
    try:
        from docx import Document
        doc = Document(file_path)
        content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        
        return "\n".join(content)
    except ImportError:
        return "Error: python-docx no está instalado"
    except Exception as e:
        return f"Error al leer {file_path}: {str(e)}"

def extract_pptx_content(file_path):
    """Extraer contenido de archivos .pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"\n--- SLIDE {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
        
        return "\n".join(content)
    except ImportError:
        return "Error: python-pptx no está instalado"
    except Exception as e:
        return f"Error al leer {file_path}: {str(e)}"

def extract_xlsx_content(file_path):
    """Extraer contenido de archivos .xlsx"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path)
        content = []
        
        for sheet_name in wb.sheetnames:
            content.append(f"\n--- HOJA: {sheet_name} ---")
            ws = wb[sheet_name]
            
            for row in ws.iter_rows(values_only=True):
                row_data = []
                for cell in row:
                    if cell is not None:
                        row_data.append(str(cell))
                if row_data:
                    content.append(" | ".join(row_data))
        
        return "\n".join(content)
    except ImportError:
        return "Error: openpyxl no está instalado"
    except Exception as e:
        return f"Error al leer {file_path}: {str(e)}"

def main():
    """Función principal"""
    docs_dir = Path("Documentacion/FASE1/Evidencias grupales")
    
    if not docs_dir.exists():
        print(f"Directorio no encontrado: {docs_dir}")
        return
    
    output_dir = Path("Documentacion/EXTRACTED")
    output_dir.mkdir(exist_ok=True)
    
    # Procesar archivos
    for file_path in docs_dir.iterdir():
        if file_path.is_file():
            print(f"Procesando: {file_path.name}")
            
            if file_path.suffix.lower() == '.docx':
                content = extract_docx_content(file_path)
                output_file = output_dir / f"{file_path.stem}.txt"
            elif file_path.suffix.lower() == '.pptx':
                content = extract_pptx_content(file_path)
                output_file = output_dir / f"{file_path.stem}.txt"
            elif file_path.suffix.lower() == '.xlsx':
                content = extract_xlsx_content(file_path)
                output_file = output_dir / f"{file_path.stem}.txt"
            else:
                print(f"Formato no soportado: {file_path.suffix}")
                continue
            
            # Guardar contenido extraído
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"CONTENIDO EXTRAÍDO DE: {file_path.name}\n")
                f.write("=" * 50 + "\n\n")
                f.write(content)
            
            print(f"Contenido guardado en: {output_file}")

if __name__ == "__main__":
    main()
